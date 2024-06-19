from pptx import Presentation
from pptx.util import Inches

# Создаем презентацию
prs = Presentation()

# Заголовок презентации
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Исследование по созданию CLI-утилиты на Python"
subtitle.text = "Доклад по результатам исследования"

# Слайд с описанием исследования
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Описание исследования"
content.text = (
    "Это исследование было направлено на создание CLI-утилиты на Python, "
    "включающей в себя несколько полезных команд, таких как арифметические операции, "
    "манипуляции с текстом, чтение файлов, а также использование API для генерации текста "
    "с помощью OpenAI GPT-4. В процессе исследования также было рассмотрено создание "
    "графического интерфейса для CLI-приложения и решения проблем с декодированием Unicode."
)

# Слайд с основными шагами и решениями
steps = [
    {
        "title": "Создание базового CLI с использованием Click",
        "content": (
            "Были созданы команды для выполнения различных операций (сложение, умножение, приветствие)."
            "\n\nПример команды:\n"
            "@cli.command()\n"
            "@click.argument('a', type=int)\n"
            "@click.argument('b', type=int)\n"
            "def add(a, b):\n"
            "    result = a + b\n"
            "    click.echo(f'The sum of {a} and {b} is {result}')"
        )
    },
    {
        "title": "Добавление новых команд",
        "content": (
            "Были добавлены команды для манипуляции с текстом и чтения файлов.\n\n"
            "Например, команда для изменения регистра текста:\n"
            "@cli.command()\n"
            "@click.option('--uppercase', '-u', is_flag=True, help='Convert text to uppercase.')\n"
            "@click.argument('text')\n"
            "def manipulate_text(text, uppercase):\n"
            "    if uppercase:\n"
            "        result = text.upper()\n"
            "    else:\n"
            "        result = text.lower()\n"
            "    click.echo(result)"
        )
    },
    {
        "title": "Решение проблемы с установкой пакета и импортом",
        "content": (
            "Было установлено правильное значение PYTHONPATH.\n\n"
            "export PYTHONPATH=/home/user/projects/cli_tool:$PYTHONPATH"
        )
    },
    {
        "title": "Создание графического интерфейса с помощью библиотеки Urwid",
        "content": (
            "Были реализованы основные функции для навигации по файловой системе.\n\n"
            "Пример использования Urwid для создания текстового интерфейса:\n"
            "import urwid\n\n"
            "class FileManagerApp:\n"
            "    def __init__(self):\n"
            "        self.widget = urwid.Text('File Manager')\n"
            "        self.loop = urwid.MainLoop(self.widget)\n"
            "        self.loop.run()\n\n"
            "if __name__ == '__main__':\n"
            "    FileManagerApp()"
        )
    },
    {
        "title": "Работа с OpenAI API для генерации текста",
        "content": (
            "Был добавлен скрипт для использования API GPT-4 для генерации текста.\n\n"
            "Решение проблем с декодированием Unicode:\n"
            "import requests\n"
            "import os\n"
            "import sys\n"
            "from dotenv import load_dotenv\n\n"
            "# Загрузка переменных окружения из файла .env\n"
            "load_dotenv()\n\n"
            "OPENAI_API_KEY = os.getenv('API_KEY')\n\n"
            "def query_gpt4o(prompt):\n"
            "    url = 'https://api.openai.com/v1/chat/completions'\n"
            "    headers = {\n"
            "        'Content-Type': 'application/json',\n"
            "        'Authorization': f'Bearer ' + OPENAI_API_KEY,\n"
            "    }\n"
            "    data = {\n"
            "        'model': 'gpt-4',\n"
            "        'messages': [{'role': 'user', 'content': prompt}],\n"
            "        'max_tokens': 50,\n"
            "    }\n"
            "    response = requests.post(url, headers=headers, json=data)\n"
            "    response_data = response.json()\n"
            "    if 'error' in response_data:\n"
            "        raise ValueError(f'Ошибка при запросе к OpenAI: {response_data['error']['message']}')\n"
            "    return response_data['choices'][0]['message']['content']\n\n"
            "if __name__ == '__main__':\n"
            "    try:\n"
            "        prompt = input('Введите текст для генерации: ').strip()\n"
            "    except UnicodeDecodeError as e:\n"
            "        print(f'Ошибка декодирования Unicode: {e}. Убедитесь, что вводите текст на русском языке в UTF-8.')\n"
            "        sys.exit(1)\n\n"
            "    try:\n"
            "        generated_text = query_gpt4o(prompt)\n"
            "        print('Сгенерированный текст:')\n"
            "        print(generated_text)\n"
            "    except ValueError as e:\n"
            "        print(e)\n"
            "    except Exception as e:\n"
            "        print(f'Произошла неизвестная ошибка: {e}')"
        )
    },
    {
        "title": "Установка корректной локали для работы с Unicode",
        "content": (
            "Было установлено значение локали для терминала:\n\n"
            "export LANG=en_US.UTF-8\n"
            "export LC_ALL=en_US.UTF-8"
        )
    }
]

for step in steps:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = step["title"]
    content.text = step["content"]

# Сохранение презентации
prs.save('CLI_Tool_Research_Report.pptx')

print("Презентация успешно создана: CLI_Tool_Research_Report.pptx")
