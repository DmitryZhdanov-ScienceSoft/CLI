from pptx import Presentation
from pptx.util import Inches
import os

# Проверим, существует ли директория
output_dir = "./"
output_file = "Fire_Library_Research_Partial.pptx"
output_path = os.path.join(output_dir, output_file)

# Создаем презентацию
prs = Presentation()

# Заголовочный слайд
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "Исследование библиотеки Fire"
subtitle.text = "Автоматическое создание CLI из объектов Python"

# Слайд с основными преимуществами
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide.shapes.title, slide.placeholders[1].text_frame
title.text = "Преимущества библиотеки Fire"

content.text = "1. Простота использования\n2. Гибкость\n3. Автоматизация"
p = content.add_paragraph()
p.text = "Простота использования: Минимум кода для создания CLI"
p = content.add_paragraph()
p.text = "Гибкость: Поддержка функций, классов и модулей"
p = content.add_paragraph()
p.text = "Автоматизация: Автоматическая генерация документации"

# Сохранение презентации
prs.save(output_path)

print(f"Презентация сохранена как {output_path}")
